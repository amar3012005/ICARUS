#!/usr/bin/env python3
"""
Extract text from a folder of real documents (PDF / docx / pptx / xlsx), chunk it, embed each
chunk with bge-m3 via the blaiq LiteLLM gateway (the engine-native 1024-dim model), and write
a corpus the mneme `.mseg` engine can ingest + a recall query set.

Outputs (bench/solvis/):
  chunks.jsonl        one line per chunk: {"id","source","chunk_idx","text"}
  corpus_f32.bin      N * 1024 float32 LE, row-major (aligned with chunks.jsonl order)
  queries_f32.bin     Q * 1024 float32 LE   (only if --query given)
  queries.txt         the query strings, one per line (only if --query given)
  meta.json           {dim, n_chunks, n_docs, model, endpoint}

Usage:
  PYTHONNOUSERSITE=1 PYTHONPATH="" bench/.venv/bin/python bench/ingest_solvis.py \
      <docs_dir> [--query "q1" --query "q2" ...]
"""
from __future__ import annotations

import argparse
import json
import os
import sys
import time
from concurrent.futures import ThreadPoolExecutor
from pathlib import Path

import numpy as np
import requests

DIM = 1024
MODEL = os.environ.get("LITELLM_EMBED_MODEL", "bge-m3")
BASE_URL = os.environ.get("LITELLM_BASE_URL", "https://api.blaiq.ai/v1").rstrip("/")
API_KEY = os.environ.get("LITELLM_API_KEY") or os.environ.get("OPENAI_API_KEY") or ""
BATCH, WORKERS, RETRIES, TIMEOUT = 64, 4, 4, 90
CHUNK_CHARS, OVERLAP, MIN_CHARS = 1500, 200, 60
OUT = Path(__file__).resolve().parent / "solvis"


# --- text extraction per file type ---------------------------------------------
def extract_pdf(p: Path) -> str:
    import fitz

    parts = []
    with fitz.open(p) as doc:
        for page in doc:
            parts.append(page.get_text("text"))
    return "\n".join(parts)


def extract_docx(p: Path) -> str:
    import docx

    d = docx.Document(str(p))
    parts = [para.text for para in d.paragraphs]
    for table in d.tables:
        for row in table.rows:
            parts.append(" | ".join(c.text for c in row.cells))
    return "\n".join(parts)


def extract_pptx(p: Path) -> str:
    from pptx import Presentation

    prs = Presentation(str(p))
    parts = []
    for i, slide in enumerate(prs.slides):
        parts.append(f"[slide {i + 1}]")
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    parts.append("".join(run.text for run in para.runs))
    return "\n".join(parts)


def extract_xlsx(p: Path) -> str:
    import openpyxl

    wb = openpyxl.load_workbook(str(p), read_only=True, data_only=True)
    parts = []
    for ws in wb.worksheets:
        parts.append(f"[sheet {ws.title}]")
        for row in ws.iter_rows(values_only=True):
            cells = [str(c) for c in row if c is not None]
            if cells:
                parts.append(" | ".join(cells))
    return "\n".join(parts)


EXTRACTORS = {
    ".pdf": extract_pdf,
    ".docx": extract_docx,
    ".pptx": extract_pptx,
    ".xlsx": extract_xlsx,
}


def chunk_text(text: str) -> list[str]:
    text = " ".join(text.split())  # collapse whitespace
    if len(text) < MIN_CHARS:
        return []
    out, i = [], 0
    while i < len(text):
        out.append(text[i : i + CHUNK_CHARS])
        i += CHUNK_CHARS - OVERLAP
    return out


# --- embedding (blaiq bge-m3) ---------------------------------------------------
def embed_batch(texts: list[str]) -> list[list[float]]:
    url = f"{BASE_URL}/embeddings"
    headers = {"authorization": f"Bearer {API_KEY}", "content-type": "application/json"}
    last = None
    for attempt in range(RETRIES):
        try:
            r = requests.post(
                url, headers=headers, json={"model": MODEL, "input": texts}, timeout=TIMEOUT
            )
            if r.status_code == 200:
                data = sorted(r.json()["data"], key=lambda d: d["index"])
                return [d["embedding"] for d in data]
            last = f"HTTP {r.status_code}: {r.text[:160]}"
        except (requests.RequestException, ValueError, KeyError) as e:
            last = repr(e)
        time.sleep(min(2**attempt, 8))
    raise RuntimeError(f"embed failed: {last}")


def embed_all(texts: list[str]) -> np.ndarray:
    batches = [texts[i : i + BATCH] for i in range(0, len(texts), BATCH)]
    results: list[list[list[float]] | None] = [None] * len(batches)
    with ThreadPoolExecutor(max_workers=WORKERS) as ex:
        futs = {ex.submit(embed_batch, b): bi for bi, b in enumerate(batches)}
        done = 0
        for fut in list(futs):
            results[futs[fut]] = fut.result()
            done += 1
            if done % 5 == 0 or done == len(batches):
                print(f"  embedded {done}/{len(batches)} batches", flush=True)
    vecs = np.zeros((len(texts), DIM), dtype=np.float32)
    w = 0
    for b in results:
        for emb in b:  # type: ignore
            v = np.asarray(emb, dtype=np.float32)
            n = np.linalg.norm(v)
            vecs[w] = v / n if n > 0 else v
            w += 1
    return vecs


def write_f32(path: Path, arr: np.ndarray) -> None:
    with open(path, "wb") as fh:
        fh.write(arr.astype("<f4", copy=False).tobytes(order="C"))


def main() -> int:
    ap = argparse.ArgumentParser()
    ap.add_argument("docs_dir")
    ap.add_argument("--query", action="append", default=[])
    ap.add_argument(
        "--queries-only",
        action="store_true",
        help="reuse existing corpus/chunks; only (re)embed the queries (fast iteration)",
    )
    args = ap.parse_args()
    if not API_KEY:
        print("FATAL: LITELLM_API_KEY not set", file=sys.stderr)
        return 1
    docs_dir = Path(args.docs_dir)
    OUT.mkdir(parents=True, exist_ok=True)

    # fast path: corpus already embedded — just embed the new queries.
    if args.queries_only:
        if not (OUT / "corpus_f32.bin").exists():
            print("FATAL: no existing corpus; run a full ingest first", file=sys.stderr)
            return 1
        if not args.query:
            print("FATAL: --queries-only needs at least one --query", file=sys.stderr)
            return 1
        qv = embed_all(args.query)
        write_f32(OUT / "queries_f32.bin", qv)
        (OUT / "queries.txt").write_text("\n".join(args.query) + "\n")
        print(f"embedded {len(args.query)} queries (reused existing corpus)", flush=True)
        return 0

    # 1. extract + chunk
    chunks: list[dict] = []
    n_docs = 0
    for p in sorted(docs_dir.iterdir()):
        if not p.is_file() or p.suffix.lower() not in EXTRACTORS:
            continue
        try:
            text = EXTRACTORS[p.suffix.lower()](p)
        except Exception as e:  # noqa: BLE001 - report + skip a bad file, don't abort
            print(f"  WARN extract {p.name}: {e}", file=sys.stderr)
            continue
        cs = chunk_text(text)
        if not cs:
            print(f"  (no text) {p.name}", flush=True)
            continue
        n_docs += 1
        for ci, c in enumerate(cs):
            chunks.append(
                {"id": len(chunks), "source": p.name, "chunk_idx": ci, "text": c}
            )
        print(f"  {p.name}: {len(cs)} chunks", flush=True)

    if not chunks:
        print("FATAL: no extractable text found", file=sys.stderr)
        return 1
    print(f"extracted {len(chunks)} chunks from {n_docs} docs. embedding {MODEL}...", flush=True)

    # 2. embed corpus
    vecs = embed_all([c["text"] for c in chunks])
    with open(OUT / "chunks.jsonl", "w") as fh:
        for c in chunks:
            fh.write(json.dumps(c, ensure_ascii=False) + "\n")
    write_f32(OUT / "corpus_f32.bin", vecs)

    # records.bin: length-prefixed "source\tchunk_idx\ttext" blobs, aligned with corpus rows.
    # The Rust demo reads this (u32 LE len + utf8 bytes) — no JSON parser needed on that side.
    with open(OUT / "records.bin", "wb") as fh:
        for c in chunks:
            blob = f"{c['source']}\t{c['chunk_idx']}\t{c['text']}".encode("utf-8")
            fh.write(len(blob).to_bytes(4, "little"))
            fh.write(blob)

    # 3. embed queries (optional)
    if args.query:
        qv = embed_all(args.query)
        write_f32(OUT / "queries_f32.bin", qv)
        (OUT / "queries.txt").write_text("\n".join(args.query) + "\n")

    (OUT / "meta.json").write_text(
        json.dumps(
            {
                "dim": DIM,
                "n_chunks": len(chunks),
                "n_docs": n_docs,
                "n_queries": len(args.query),
                "model": MODEL,
                "endpoint": f"{BASE_URL}/embeddings",
            },
            indent=2,
        )
    )
    print(
        f"done. {len(chunks)} chunks / {n_docs} docs -> {OUT}/ "
        f"(corpus_f32.bin, chunks.jsonl{', queries_f32.bin' if args.query else ''})",
        flush=True,
    )
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
